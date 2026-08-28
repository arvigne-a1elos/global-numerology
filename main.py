# -*- coding: utf-8 -*-
# main.py - A1ELOS Global Numerology API
# VERSÃO CONSOLIDADA CORRIGIDA - TODAS AS ROTAS DOS 23 PRODUTOS
import stripe
from produtos.mapa import reduzir, calc_mapa, calc_grid, analisar_express, analisar_completo, analisar_vida
from produtos.nome import analisar_nome
from produtos.urna import validar_nomes_urna
from produtos.eleitoral import gerar_numeros
from produtos.imovel import analisar_imovel
from produtos.calendario import analisar_calendario
from produtos.casal import analisar_casal
from produtos.familia import analisar_familia
from produtos.coletivo import desconto_bc
import os, json, uuid, logging, secrets, string, base64, traceback
from datetime import date, datetime
from typing import Optional
from fastapi import FastAPI, HTTPException, Request
from fastapi.responses import FileResponse, HTMLResponse, RedirectResponse
from fastapi.staticfiles import StaticFiles
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from sqlalchemy import create_engine, Column, String, Integer, Float, DateTime
from sqlalchemy.orm import sessionmaker, declarative_base
from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import ParagraphStyle
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT
import qrcode
import dateutil.parser as dp
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from apresentacao_textos import APRESENTACAO_TEXTOS
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)
STRIPE_KEY = os.getenv("STRIPE_SECRET_KEY", "")
STRIPE_PUB = os.getenv("STRIPE_PUBLISHABLE_KEY", "")
FROM_EMAIL = os.getenv("FROM_EMAIL", "noreply@a1elos.com.br")
FROM_NAME = "A1ELOS Numerologia Global"
BASE_URL = os.getenv("BASE_URL", os.getenv("SITE_URL", "https://global-numerology.onrender.com"))
DB_URL = os.getenv("DATABASE_URL", "sqlite:///./numerologia.db")
ADMIN_EMAIL = "arvigne@gmail.com"
if STRIPE_KEY:
    stripe.api_key = STRIPE_KEY
# ===== BANCO DE DADOS =====
engine_kwargs = {}
if DB_URL.startswith("sqlite"):
    engine_kwargs["connect_args"] = {"check_same_thread": False}
engine = create_engine(DB_URL, **engine_kwargs)
Base = declarative_base()
SessionLocal = sessionmaker(bind=engine)
class Calc(Base):
    __tablename__ = "calculations"
    id = Column(String, primary_key=True)
    name = Column(String)
    birth_date = Column(String)
    email = Column(String, nullable=True)
    life_path = Column(Integer)
    expression = Column(Integer)
    soul_urge = Column(Integer)
    personality = Column(Integer)
    destiny = Column(Integer)
    created_at = Column(DateTime, default=datetime.utcnow)
class Order(Base):
    __tablename__ = "orders"
    id = Column(String, primary_key=True)
    email = Column(String)
    product = Column(String)
    price = Column(Float)
    status = Column(String, default="pending")
    payment_id = Column(String, nullable=True)
    created_at = Column(DateTime, default=datetime.utcnow)
try:
    Base.metadata.create_all(bind=engine)
except Exception as e:
    logger.error(f"DB init adiado: {e}")
# ===== APP =====
app = FastAPI(title="Global Numerology")
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_credentials=True,
                   allow_methods=["*"], allow_headers=["*"])
STATIC_DIR = os.path.join(os.path.dirname(__file__), "static")
if os.path.isdir(STATIC_DIR):
    app.mount("/static", StaticFiles(directory=STATIC_DIR), name="static")
# ===== 12 IDIOMAS E MOEDAS =====
IDIOMAS = ["pt", "en", "es", "it", "fr", "de", "ja", "zh", "ru", "hi", "he", "ar"]
MOEDA = {
    "pt": "brl", "en": "usd", "es": "eur", "it": "eur", "fr": "eur", "de": "eur",
    "ja": "jpy", "zh": "cny", "ru": "rub", "hi": "inr", "he": "ils", "ar": "sar"
}
SIMBOLO = {
    "pt": "R$", "en": "US$", "es": "€", "it": "€", "fr": "€", "de": "€",
    "ja": "¥", "zh": "¥", "ru": "₽", "hi": "₹", "he": "₪", "ar": "﷼"
}
# ===== FAIXAS DE PREÇO (23 produtos) =====
PRODUTO_FAIXA = {
    "express": 0, "vida": 0, "completo": 1, "ia": 1,
    "urna": 2, "eleitoral": 2, "imovel": 2, "calendario": 2,
    "artistico": 3, "bebe": 3, "assinatura": 3,
    "negocio": 4, "casal": 4, "familia": 5,
    "nome_pet": 0, "nickname": 0, "nome_dominio": 0, "nome_canal": 0,
    "nome_equipe": 0, "nome_ong": 0, "nome_projeto": 0, "nome_evento": 0
}
VALORES = {
    "pt": [800, 1700, 2600, 3500, 4400, 9800],
    "en": [150, 350, 500, 700, 900, 2000],
    "es": [150, 350, 500, 700, 900, 2000],
    "it": [150, 350, 500, 700, 900, 2000],
    "fr": [150, 350, 500, 700, 900, 2000],
    "de": [150, 350, 500, 700, 900, 2000],
    "ja": [250, 550, 800, 1100, 1400, 3200],
    "zh": [1200, 2500, 4000, 5500, 7000, 16000],
    "ru": [13000, 28000, 43000, 58000, 73000, 160000],
    "hi": [12000, 28000, 43000, 58000, 73000, 160000],
    "he": [500, 1300, 1900, 2600, 3300, 7300],
    "ar": [600, 1300, 1900, 2600, 3300, 7300]
}
def preco_local(produto, lang):
    return VALORES[lang][PRODUTO_FAIXA[produto]]
# ===== NOMES DOS 23 PRODUTOS (12 IDIOMAS) =====
PRODUTOS = {
    "pt": {"express": "Mapa Express", "vida": "Qual Vida/Ano", "completo": "Mapa Completo",
        "ia": "Pesquisa IA de Nomes", "urna": "Validacao Nome de Urna", "eleitoral": "Numero Eleitoral",
        "imovel": "Numero do Imovel", "calendario": "Calendario Mensal Energetico",
        "artistico": "Validacao Nome Artistico", "bebe": "Planejamento Nome de Bebe",
        "assinatura": "Validacao de Assinaturas", "negocio": "Nome para Negocio/Produto",
        "casal": "Mapa do Casal", "familia": "Mapa Familia Premium", "coletivo": "Bonus Coletivo/Empresarial",
        "nome_pet": "Nome do Pet", "nickname": "Nickname Digital", "nome_dominio": "Nome do Dominio",
        "nome_canal": "Nome do Canal", "nome_equipe": "Nome da Equipe", "nome_ong": "Nome de ONG, Associacao, Instituto ou Fundacao",
        "nome_projeto": "Nome do Projeto", "nome_evento": "Nome do Evento"},
    "en": {"express": "Express Map", "vida": "Life Phase & Year", "completo": "Complete Map",
        "ia": "AI Name Search", "urna": "Ballot Name Validation", "eleitoral": "Electoral Number",
        "imovel": "Property Number", "calendario": "Monthly Energy Calendar",
        "artistico": "Artistic Name Validation", "bebe": "Baby Name Planning",
        "assinatura": "Signature Validation", "negocio": "Business & Product Name",
        "casal": "Couple Map", "familia": "Premium Family Map", "coletivo": "Corporate Bonus",
        "nome_pet": "Pet Name", "nickname": "Digital Nickname", "nome_dominio": "Domain Name",
        "nome_canal": "Channel Name", "nome_equipe": "Team Name", "nome_ong": "NGO, Association, Institute or Foundation Name",
        "nome_projeto": "Project Name", "nome_evento": "Event Name"},
    "es": {"express": "Mapa Exprés", "vida": "Ciclo de Vida y Año", "completo": "Mapa Completo",
        "ia": "Búsqueda IA de Nombres", "urna": "Validación Nombre de Urna", "eleitoral": "Número Electoral",
        "imovel": "Número de la Propiedad", "calendario": "Calendario Mensual Energético",
        "artistico": "Validación Nombre Artístico", "bebe": "Planificación Nombre de Bebé",
        "assinatura": "Validación de Firmas", "negocio": "Nombre para Negocio/Producto",
        "casal": "Mapa de Pareja", "familia": "Mapa Familiar Premium", "coletivo": "Bono Corporativo",
        "nome_pet": "Nombre de la Mascota", "nickname": "Apodo Digital", "nome_dominio": "Nombre de Dominio",
        "nome_canal": "Nombre del Canal", "nome_equipe": "Nombre del Equipo", "nome_ong": "Nombre de ONG, Asociacion, Instituto o Fundacion",
        "nome_projeto": "Nombre del Proyecto", "nome_evento": "Nombre del Evento"},
    "it": {"express": "Mappa Espressa", "vida": "Fase di Vita e Anno", "completo": "Mappa Completa",
        "ia": "Ricerca IA Nomi", "urna": "Validazione Nome della Scheda", "eleitoral": "Numero Elettorale",
        "imovel": "Numero dell'Immobile", "calendario": "Calendario Mensile Energetico",
        "artistico": "Validazione Nome d'Arte", "bebe": "Pianificazione Nome del Bambino",
        "assinatura": "Validazione delle Firme", "negocio": "Nome per Business/Prodotto",
        "casal": "Mappa di Coppia", "familia": "Mappa Famiglia Premium", "coletivo": "Bonus Aziendale",
        "nome_pet": "Nome dell'Animale", "nickname": "Nickname Digitale", "nome_dominio": "Nome del Dominio",
        "nome_canal": "Nome del Canale", "nome_equipe": "Nome della Squadra", "nome_ong": "Nome di ONG, Associazione, Istituto o Fondazione",
        "nome_projeto": "Nome del Progetto", "nome_evento": "Nome dell'Evento"},
    "fr": {"express": "Carte Express", "vida": "Phase de Vie et Année", "completo": "Carte Complète",
        "ia": "Recherche IA de Noms", "urna": "Validation Nom du Bulletin", "eleitoral": "Numéro Électoral",
        "imovel": "Numéro du Bien", "calendario": "Calendrier Mensuel Énergétique",
        "artistico": "Validation Nom de Scène", "bebe": "Planification Prénom de Bébé",
        "assinatura": "Validation des Signatures", "negocio": "Nom pour Entreprise/Produit",
        "casal": "Carte du Couple", "familia": "Carte Famille Premium", "coletivo": "Bonus d'Entreprise",
        "nome_pet": "Nom de l'Animal", "nickname": "Pseudo Numerique", "nome_dominio": "Nom de Domaine",
        "nome_canal": "Nom de la Chaine", "nome_equipe": "Nom de l'Equipe", "nome_ong": "Nom d'ONG, Association, Institut ou Fondation",
        "nome_projeto": "Nom du Projet", "nome_evento": "Nom de l'Evenement"},
    "de": {"express": "Express-Karte", "vida": "Lebensphase & Jahr", "completo": "Vollständige Karte",
        "ia": "KI-Namenssuche", "urna": "Stimmzettelname-Validierung", "eleitoral": "Wahlnummer",
        "imovel": "Immobiliennummer", "calendario": "Monatlicher Energiekalender",
        "artistico": "Künstlername-Validierung", "bebe": "Babynamen-Planung",
        "assinatura": "Unterschrifts-Validierung", "negocio": "Name für Unternehmen/Produkt",
        "casal": "Paar-Karte", "familia": "Premium-Familien-Karte", "coletivo": "Unternehmensbonus",
        "nome_pet": "Haustiername", "nickname": "Digitaler Spitzname", "nome_dominio": "Domainname",
        "nome_canal": "Kanalname", "nome_equipe": "Teamname", "nome_ong": "Name von NGO, Verein, Institut oder Stiftung",
        "nome_projeto": "Projektname", "nome_evento": "Veranstaltungsname"},
    "ja": {"express": "エクスプレスマップ", "vida": "ライフステージと年", "completo": "完全マップ",
        "ia": "AI名前検索", "urna": "投票用紙名の検証", "eleitoral": "選挙番号",
        "imovel": "不動産番号", "calendario": "月間エネルギーカレンダー",
        "artistico": "芸名の検証", "bebe": "赤ちゃんの名前計画",
        "assinatura": "署名の検証", "negocio": "ビジネス・商品名",
        "casal": "カップルマップ", "familia": "プレミアム家族マップ", "coletivo": "法人ボーナス",
        "nome_pet": "ペットの名前", "nickname": "デジタルニックネーム", "nome_dominio": "ドメイン名",
        "nome_canal": "チャンネル名", "nome_equipe": "チーム名", "nome_ong": "NGO・協会・研究所・財団の名前",
        "nome_projeto": "プロジェクト名", "nome_evento": "イベント名"},
    "zh": {"express": "快速地图", "vida": "生命阶段与年份", "completo": "完整地图",
        "ia": "AI名字搜索", "urna": "选票名称验证", "eleitoral": "选举号码",
        "imovel": "房产号码", "calendario": "每月能量日历",
        "artistico": "艺名验证", "bebe": "宝宝取名规划",
        "assinatura": "签名验证", "negocio": "企业/产品名称",
        "casal": "情侣地图", "familia": "高级家庭地图", "coletivo": "企业奖励",
        "nome_pet": "宠物名字", "nickname": "数字昵称", "nome_dominio": "域名",
        "nome_canal": "频道名称", "nome_equipe": "团队名称", "nome_ong": "非政府组织、协会、研究所或基金会名称",
        "nome_projeto": "项目名称", "nome_evento": "活动名称"},
    "ru": {"express": "Экспресс-карта", "vida": "Жизненный этап и год", "completo": "Полная карта",
        "ia": "ИИ-поиск имён", "urna": "Проверка названия бюллетеня", "eleitoral": "Избирательный номер",
        "imovel": "Номер недвижимости", "calendario": "Ежемесячный энергетический календарь",
        "artistico": "Проверка сценического имени", "bebe": "Планирование имени ребёнка",
        "assinatura": "Проверка подписей", "negocio": "Название для бизнеса/продукта",
        "casal": "Карта пары", "familia": "Премиальная семейная карта", "coletivo": "Корпоративный бонус",
        "nome_pet": "Имя питомца", "nickname": "Цифровой никнейм", "nome_dominio": "Имя домена",
        "nome_canal": "Название канала", "nome_equipe": "Название команды", "nome_ong": "Название НКО, ассоциации, института или фонда",
        "nome_projeto": "Название проекта", "nome_evento": "Название события"},
    "hi": {"express": "त्वरित मानचित्र", "vida": "जीवन चरण और वर्ष", "completo": "पूर्ण मानचित्र",
        "ia": "AI नाम खोज", "urna": "मतपत्र नाम सत्यापन", "eleitoral": "निर्वाचन संख्या",
        "imovel": "संपत्ति संख्या", "calendario": "मासिक ऊर्जा कैलेंडर",
        "artistico": "कलात्मक नाम सत्यापन", "bebe": "शिशु नाम योजना",
        "assinatura": "हस्ताक्षर सत्यापन", "negocio": "व्यवसाय/उत्पाद नाम",
        "casal": "युगल मानचित्र", "familia": "प्रीमियम परिवार मानचित्र", "coletivo": "कॉर्पोरेट बोनस",
        "nome_pet": "पालतू नाम", "nickname": "डिजिटल उपनाम", "nome_dominio": "डोमेन नाम",
        "nome_canal": "चैनल नाम", "nome_equipe": "टीम नाम", "nome_ong": "एनजीओ, संघ, संस्थान या फाउंडेशन का नाम",
        "nome_projeto": "परियोजना नाम", "nome_evento": "कार्यक्रम नाम"},
    "he": {"express": "מפה מהירה", "vida": "שלב חיים ושנה", "completo": "מפה מלאה",
        "ia": "חיפוש שמות AI", "urna": "אימות שם פתק", "eleitoral": "מספר בחירות",
        "imovel": "מספר נכס", "calendario": "לוח אנרגיה חודשי",
        "artistico": "אימות שם במה", "bebe": "תכנון שם לתינוק",
        "assinatura": "אימות חתימות", "negocio": "שם לעסק/מוצר",
        "nome_pet": "שם חיית המחמד", "nickname": "כינוי דיגיטלי", "nome_dominio": "שם דומיין",
        "nome_canal": "שם הערוץ", "nome_equipe": "שם הצוות", "nome_ong": "שם עמותה, ארגון, מכון או קרן",
        "nome_projeto": "שם הפרויקט", "nome_evento": "שם האירוע",
        "casal": "מפת זוג", "familia": "מפת משפחה פרימיום", "coletivo": "בונוס ארגוני"},
    "ar": {"express": "خريطة سريعة", "vida": "مرحلة الحياة والسنة", "completo": "خريطة كاملة",
        "ia": "بحث الأسماء بالذكاء الاصطناعي", "urna": "التحقق من اسم الاقتراع", "eleitoral": "الرقم الانتخابي",
        "imovel": "رقم العقار", "calendario": "التقويم الشهري للطاقة",
        "artistico": "التحقق من الاسم الفني", "bebe": "تخطيط اسم الطفل",
        "assinatura": "التحقق من التوقيعات", "negocio": "اسم للأعمال/المنتج",
        "casal": "خريطة الزوجين", "familia": "خريطة العائلة المميزة", "coletivo": "مكافأة الشركات",
        "nome_pet": "اسم الحيوان الأليف", "nickname": "اللقب الرقمي", "nome_dominio": "اسم النطاق",
        "nome_canal": "اسم القناة", "nome_equipe": "اسم الفريق", "nome_ong": "اسم منظمة أو جمعية أو معهد أو مؤسسة",
        "nome_projeto": "اسم المشروع", "nome_evento": "اسم الفعالية"}
}
# ===== PRICE IDS STRIPE (23 produtos, 12 idiomas) =====
PRICE_IDS = {
    "pt": {"express": "price_1TxocVBMLa84bVJ0EL0kb9Dn", "completo": "price_1TxohlBMLa84bVJ0jVj9307b",
           "urna": "price_1TxollBMLa84bVJ0Wk5zIak6", "eleitoral": "price_1TxopFBMLa84bVJ0jvtJExVj",
           "vida": "PRICE_ID_PT_VIDA", "ia": "PRICE_ID_PT_IA", "imovel": "PRICE_ID_PT_IMOVEL",
           "calendario": "PRICE_ID_PT_CALENDARIO", "artistico": "PRICE_ID_PT_ARTISTICO",
           "bebe": "PRICE_ID_PT_BEBE", "assinatura": "PRICE_ID_PT_ASSINATURA",
           "negocio": "PRICE_ID_PT_NEGOCIO", "casal": "PRICE_ID_PT_CASAL", "familia": "PRICE_ID_PT_FAMILIA",
           "nome_pet": "PRICE_ID_PT_NOME_PET", "nickname": "PRICE_ID_PT_NICKNAME",
           "nome_dominio": "PRICE_ID_PT_NOME_DOMINIO", "nome_canal": "PRICE_ID_PT_NOME_CANAL",
           "nome_equipe": "PRICE_ID_PT_NOME_EQUIPE", "nome_ong": "PRICE_ID_PT_NOME_ONG",
           "nome_projeto": "PRICE_ID_PT_NOME_PROJETO", "nome_evento": "PRICE_ID_PT_NOME_EVENTO"},
    "en": {"express": "price_1TxotnBMLa84bVJ00SGo4kjO", "completo": "price_1TxoxfBMLa84bVJ0VgQVddZX",
           "urna": "price_1Txp1jBMLa84bVJ06W4559rN", "eleitoral": "price_1Txp5aBMLa84bVJ0GqrvBrIk",
           "vida": "PRICE_ID_EN_VIDA", "ia": "PRICE_ID_EN_IA", "imovel": "PRICE_ID_EN_IMOVEL",
           "calendario": "PRICE_ID_EN_CALENDARIO", "artistico": "PRICE_ID_EN_ARTISTICO",
           "bebe": "PRICE_ID_EN_BEBE", "assinatura": "PRICE_ID_EN_ASSINATURA",
           "negocio": "PRICE_ID_EN_NEGOCIO", "casal": "PRICE_ID_EN_CASAL", "familia": "PRICE_ID_EN_FAMILIA",
           "nome_pet": "PRICE_ID_EN_NOME_PET", "nickname": "PRICE_ID_EN_NICKNAME",
           "nome_dominio": "PRICE_ID_EN_NOME_DOMINIO", "nome_canal": "PRICE_ID_EN_NOME_CANAL",
           "nome_equipe": "PRICE_ID_EN_NOME_EQUIPE", "nome_ong": "PRICE_ID_EN_NOME_ONG",
           "nome_projeto": "PRICE_ID_EN_NOME_PROJETO", "nome_evento": "PRICE_ID_EN_NOME_EVENTO"},
    "es": {"express": "price_1TyD2oBMLa84bVJ0HvSTMozS", "completo": "price_1TyD6NBMLa84bVJ0s5y2OtSr",
           "urna": "price_1TyDB0BMLa84bVJ0baUEGa2P", "eleitoral": "price_1TyDCsBMLa84bVJ0NRp5uOKU",
           "vida": "PRICE_ID_ES_VIDA", "ia": "PRICE_ID_ES_IA", "imovel": "PRICE_ID_ES_IMOVEL",
           "calendario": "PRICE_ID_ES_CALENDARIO", "artistico": "PRICE_ID_ES_ARTISTICO",
           "bebe": "PRICE_ID_ES_BEBE", "assinatura": "PRICE_ID_ES_ASSINATURA",
           "negocio": "PRICE_ID_ES_NEGOCIO", "casal": "PRICE_ID_ES_CASAL", "familia": "PRICE_ID_ES_FAMILIA",
           "nome_pet": "PRICE_ID_ES_NOME_PET", "nickname": "PRICE_ID_ES_NICKNAME",
           "nome_dominio": "PRICE_ID_ES_NOME_DOMINIO", "nome_canal": "PRICE_ID_ES_NOME_CANAL",
           "nome_equipe": "PRICE_ID_ES_NOME_EQUIPE", "nome_ong": "PRICE_ID_ES_NOME_ONG",
           "nome_projeto": "PRICE_ID_ES_NOME_PROJETO", "nome_evento": "PRICE_ID_ES_NOME_EVENTO"},
    "it": {"express": "PRICE_ID_IT_EXPRESS", "completo": "PRICE_ID_IT_COMPLETO",
           "urna": "PRICE_ID_IT_URNA", "eleitoral": "PRICE_ID_IT_ELEITORAL",
           "vida": "PRICE_ID_IT_VIDA", "ia": "PRICE_ID_IT_IA", "imovel": "PRICE_ID_IT_IMOVEL",
           "calendario": "PRICE_ID_IT_CALENDARIO", "artistico": "PRICE_ID_IT_ARTISTICO",
           "bebe": "PRICE_ID_IT_BEBE", "assinatura": "PRICE_ID_IT_ASSINATURA",
           "negocio": "PRICE_ID_IT_NEGOCIO", "casal": "PRICE_ID_IT_CASAL", "familia": "PRICE_ID_IT_FAMILIA",
           "nome_pet": "PRICE_ID_IT_NOME_PET", "nickname": "PRICE_ID_IT_NICKNAME",
           "nome_dominio": "PRICE_ID_IT_NOME_DOMINIO", "nome_canal": "PRICE_ID_IT_NOME_CANAL",
           "nome_equipe": "PRICE_ID_IT_NOME_EQUIPE", "nome_ong": "PRICE_ID_IT_NOME_ONG",
           "nome_projeto": "PRICE_ID_IT_NOME_PROJETO", "nome_evento": "PRICE_ID_IT_NOME_EVENTO"},
    "fr": {"express": "price_1TyDnQBMLa84bVJ0K9DBz2mk", "completo": "price_1TyDrjBMLa84bVJ0cstgcPbY",
           "urna": "price_1TyDw1BMLa84bVJ0EV0OnINW", "eleitoral": "price_1TyDxsBMLa84bVJ0n2t4jOfZ",
           "vida": "PRICE_ID_FR_VIDA", "ia": "PRICE_ID_FR_IA", "imovel": "PRICE_ID_FR_IMOVEL",
           "calendario": "PRICE_ID_FR_CALENDARIO", "artistico": "PRICE_ID_FR_ARTISTICO",
           "bebe": "PRICE_ID_FR_BEBE", "assinatura": "PRICE_ID_FR_ASSINATURA",
           "negocio": "PRICE_ID_FR_NEGOCIO", "casal": "PRICE_ID_FR_CASAL", "familia": "PRICE_ID_FR_FAMILIA",
           "nome_pet": "PRICE_ID_FR_NOME_PET", "nickname": "PRICE_ID_FR_NICKNAME",
           "nome_dominio": "PRICE_ID_FR_NOME_DOMINIO", "nome_canal": "PRICE_ID_FR_NOME_CANAL",
           "nome_equipe": "PRICE_ID_FR_NOME_EQUIPE", "nome_ong": "PRICE_ID_FR_NOME_ONG",
           "nome_projeto": "PRICE_ID_FR_NOME_PROJETO", "nome_evento": "PRICE_ID_FR_NOME_EVENTO"},
    "de": {"express": "price_1TyFJaBMLa84bVJ0BDPNQUjz", "completo": "price_1TyFLKBMLa84bVJ0RT0bkKpW",
           "urna": "price_1TyFO2BMLa84bVJ0FIoh7co1", "eleitoral": "price_1TyFTxBMLa84bVJ0qw6LQvVI",
           "vida": "PRICE_ID_DE_VIDA", "ia": "PRICE_ID_DE_IA", "imovel": "PRICE_ID_DE_IMOVEL",
           "calendario": "PRICE_ID_DE_CALENDARIO", "artistico": "PRICE_ID_DE_ARTISTICO",
           "bebe": "PRICE_ID_DE_BEBE", "assinatura": "PRICE_ID_DE_ASSINATURA",
           "negocio": "PRICE_ID_DE_NEGOCIO", "casal": "PRICE_ID_DE_CASAL", "familia": "PRICE_ID_DE_FAMILIA",
           "nome_pet": "PRICE_ID_DE_NOME_PET", "nickname": "PRICE_ID_DE_NICKNAME",
           "nome_dominio": "PRICE_ID_DE_NOME_DOMINIO", "nome_canal": "PRICE_ID_DE_NOME_CANAL",
           "nome_equipe": "PRICE_ID_DE_NOME_EQUIPE", "nome_ong": "PRICE_ID_DE_NOME_ONG",
           "nome_projeto": "PRICE_ID_DE_NOME_PROJETO", "nome_evento": "PRICE_ID_DE_NOME_EVENTO"},
    "ja": {"express": "price_1TyJ5HBMLa84bVJ00nZLnuV1", "completo": "price_1TyJJgBMLa84bVJ0fkO5nSFT",
           "urna": "price_1TyJOzBMLa84bVJ0BAPegYVD", "eleitoral": "price_1TyJRwBMLa84bVJ0PLA1CIuH",
           "vida": "PRICE_ID_JA_VIDA", "ia": "PRICE_ID_JA_IA", "imovel": "PRICE_ID_JA_IMOVEL",
           "calendario": "PRICE_ID_JA_CALENDARIO", "artistico": "PRICE_ID_JA_ARTISTICO",
           "bebe": "PRICE_ID_JA_BEBE", "assinatura": "PRICE_ID_JA_ASSINATURA",
           "negocio": "PRICE_ID_JA_NEGOCIO", "casal": "PRICE_ID_JA_CASAL", "familia": "PRICE_ID_JA_FAMILIA",
           "nome_pet": "PRICE_ID_JA_NOME_PET", "nickname": "PRICE_ID_JA_NICKNAME",
           "nome_dominio": "PRICE_ID_JA_NOME_DOMINIO", "nome_canal": "PRICE_ID_JA_NOME_CANAL",
           "nome_equipe": "PRICE_ID_JA_NOME_EQUIPE", "nome_ong": "PRICE_ID_JA_NOME_ONG",
           "nome_projeto": "PRICE_ID_JA_NOME_PROJETO", "nome_evento": "PRICE_ID_JA_NOME_EVENTO"},
    "zh": {"express": "price_1TyKXeBMLa84bVJ07Q6w0j6G", "completo": "price_1TyKZfBMLa84bVJ0bgYSm8e2",
           "urna": "price_1TyKdWBMLa84bVJ0TIP0Knbi", "eleitoral": "price_1TyKitBMLa84bVJ0lFgyKya0",
           "vida": "PRICE_ID_ZH_VIDA", "ia": "PRICE_ID_ZH_IA", "imovel": "PRICE_ID_ZH_IMOVEL",
           "calendario": "PRICE_ID_ZH_CALENDARIO", "artistico": "PRICE_ID_ZH_ARTISTICO",
           "bebe": "PRICE_ID_ZH_BEBE", "assinatura": "PRICE_ID_ZH_ASSINATURA",
           "negocio": "PRICE_ID_ZH_NEGOCIO", "casal": "PRICE_ID_ZH_CASAL", "familia": "PRICE_ID_ZH_FAMILIA",
           "nome_pet": "PRICE_ID_ZH_NOME_PET", "nickname": "PRICE_ID_ZH_NICKNAME",
           "nome_dominio": "PRICE_ID_ZH_NOME_DOMINIO", "nome_canal": "PRICE_ID_ZH_NOME_CANAL",
           "nome_equipe": "PRICE_ID_ZH_NOME_EQUIPE", "nome_ong": "PRICE_ID_ZH_NOME_ONG",
           "nome_projeto": "PRICE_ID_ZH_NOME_PROJETO", "nome_evento": "PRICE_ID_ZH_NOME_EVENTO"},
    "ru": {"express": "PRICE_ID_RU_EXPRESS", "completo": "PRICE_ID_RU_COMPLETO",
           "urna": "PRICE_ID_RU_URNA", "eleitoral": "PRICE_ID_RU_ELEITORAL",
           "vida": "PRICE_ID_RU_VIDA", "ia": "PRICE_ID_RU_IA", "imovel": "PRICE_ID_RU_IMOVEL",
           "calendario": "PRICE_ID_RU_CALENDARIO", "artistico": "PRICE_ID_RU_ARTISTICO",
           "bebe": "PRICE_ID_RU_BEBE", "assinatura": "PRICE_ID_RU_ASSINATURA",
           "negocio": "PRICE_ID_RU_NEGOCIO", "casal": "PRICE_ID_RU_CASAL", "familia": "PRICE_ID_RU_FAMILIA",
           "nome_pet": "PRICE_ID_RU_NOME_PET", "nickname": "PRICE_ID_RU_NICKNAME",
           "nome_dominio": "PRICE_ID_RU_NOME_DOMINIO", "nome_canal": "PRICE_ID_RU_NOME_CANAL",
           "nome_equipe": "PRICE_ID_RU_NOME_EQUIPE", "nome_ong": "PRICE_ID_RU_NOME_ONG",
           "nome_projeto": "PRICE_ID_RU_NOME_PROJETO", "nome_evento": "PRICE_ID_RU_NOME_EVENTO"},
    "hi": {"express": "PRICE_ID_HI_EXPRESS", "completo": "PRICE_ID_HI_COMPLETO",
           "urna": "PRICE_ID_HI_URNA", "eleitoral": "PRICE_ID_HI_ELEITORAL",
           "vida": "PRICE_ID_HI_VIDA", "ia": "PRICE_ID_HI_IA", "imovel": "PRICE_ID_HI_IMOVEL",
           "calendario": "PRICE_ID_HI_CALENDARIO", "artistico": "PRICE_ID_HI_ARTISTICO",
           "bebe": "PRICE_ID_HI_BEBE", "assinatura": "PRICE_ID_HI_ASSINATURA",
           "negocio": "PRICE_ID_HI_NEGOCIO", "casal": "PRICE_ID_HI_CASAL", "familia": "PRICE_ID_HI_FAMILIA",
           "nome_pet": "PRICE_ID_HI_NOME_PET", "nickname": "PRICE_ID_HI_NICKNAME",
           "nome_dominio": "PRICE_ID_HI_NOME_DOMINIO", "nome_canal": "PRICE_ID_HI_NOME_CANAL",
           "nome_equipe": "PRICE_ID_HI_NOME_EQUIPE", "nome_ong": "PRICE_ID_HI_NOME_ONG",
           "nome_projeto": "PRICE_ID_HI_NOME_PROJETO", "nome_evento": "PRICE_ID_HI_NOME_EVENTO"},
    "he": {"express": "price_1TyIKeBMLa84bVJ0W02dbXOt", "completo": "price_1TyIO0BMLa84bVJ08P0j9THk",
           "urna": "price_1TyIPbBMLa84bVJ08GnGksRk", "eleitoral": "price_1TyISQBMLa84bVJ0sb7xjIyV",
           "vida": "PRICE_ID_HE_VIDA", "ia": "PRICE_ID_HE_IA", "imovel": "PRICE_ID_HE_IMOVEL",
           "calendario": "PRICE_ID_HE_CALENDARIO", "artistico": "PRICE_ID_HE_ARTISTICO",
           "bebe": "PRICE_ID_HE_BEBE", "assinatura": "PRICE_ID_HE_ASSINATURA",
           "negocio": "PRICE_ID_HE_NEGOCIO", "casal": "PRICE_ID_HE_CASAL", "familia": "PRICE_ID_HE_FAMILIA",
           "nome_pet": "PRICE_ID_HE_NOME_PET", "nickname": "PRICE_ID_HE_NICKNAME",
           "nome_dominio": "PRICE_ID_HE_NOME_DOMINIO", "nome_canal": "PRICE_ID_HE_NOME_CANAL",
           "nome_equipe": "PRICE_ID_HE_NOME_EQUIPE", "nome_ong": "PRICE_ID_HE_NOME_ONG",
           "nome_projeto": "PRICE_ID_HE_NOME_PROJETO", "nome_evento": "PRICE_ID_HE_NOME_EVENTO"},
    "ar": {"express": "PRICE_ID_AR_EXPRESS", "completo": "PRICE_ID_AR_COMPLETO",
           "urna": "PRICE_ID_AR_URNA", "eleitoral": "PRICE_ID_AR_ELEITORAL",
           "vida": "PRICE_ID_AR_VIDA", "ia": "PRICE_ID_AR_IA", "imovel": "PRICE_ID_AR_IMOVEL",
           "calendario": "PRICE_ID_AR_CALENDARIO", "artistico": "PRICE_ID_AR_ARTISTICO",
           "bebe": "PRICE_ID_AR_BEBE", "assinatura": "PRICE_ID_AR_ASSINATURA",
           "negocio": "PRICE_ID_AR_NEGOCIO", "casal": "PRICE_ID_AR_CASAL", "familia": "PRICE_ID_AR_FAMILIA",
           "nome_pet": "PRICE_ID_AR_NOME_PET", "nickname": "PRICE_ID_AR_NICKNAME",
           "nome_dominio": "PRICE_ID_AR_NOME_DOMINIO", "nome_canal": "PRICE_ID_AR_NOME_CANAL",
           "nome_equipe": "PRICE_ID_AR_NOME_EQUIPE", "nome_ong": "PRICE_ID_AR_NOME_ONG",
           "nome_projeto": "PRICE_ID_AR_NOME_PROJETO", "nome_evento": "PRICE_ID_AR_NOME_EVENTO"}
}
PRODUTO_TARGET = {
    "express": "calculadora", "vida": "produtos", "completo": "calculadora",
    "ia": "produtos", "urna": "form-urna", "eleitoral": "form-eleitoral",
    "imovel": "produtos", "calendario": "produtos", "artistico": "produtos",
    "bebe": "produtos", "assinatura": "produtos", "negocio": "produtos",
    "casal": "produtos", "familia": "produtos", "coletivo": "corporativo",
    "nome_pet": "calculadora", "nickname": "calculadora", "nome_dominio": "calculadora",
    "nome_canal": "calculadora", "nome_equipe": "calculadora", "nome_ong": "calculadora",
    "nome_projeto": "calculadora", "nome_evento": "calculadora"
}

# ===== APRESENTAÇÃO DE SLIDES POR IDIOMA (PPTX dinâmico) =====
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

GOLD_HEX = RGBColor(0xB8, 0x86, 0x0B)
DARK_HEX = RGBColor(0x22, 0x22, 0x22)
GRAY_HEX = RGBColor(0x88, 0x88, 0x88)

def _slide_texto(slide, left, top, width, height, texto, tam=14,
                 negrito=False, cor=DARK_HEX, alinh=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = texto
    p.font.size = Pt(tam)
    p.font.bold = negrito
    p.font.color.rgb = cor
    p.alignment = alinh
    return box

def _slide_titulo(prs, t):
    s = prs.slides.add_slide(prs.slide_layouts[6])          # layout em branco
    _slide_texto(s, 1.0, 2.0, 8.0, 1.0, "A1ELOS GLOBAL NUMEROLOGY", 18, True, GOLD_HEX, PP_ALIGN.CENTER)
    _slide_texto(s, 1.0, 3.0, 8.0, 0.7, t["subtitulo"], 13, False, GRAY_HEX, PP_ALIGN.CENTER)
    _slide_texto(s, 1.0, 3.8, 8.0, 1.4, t["titulo"], 28, True, DARK_HEX, PP_ALIGN.CENTER)
    _slide_texto(s, 1.0, 5.6, 8.0, 0.6, t["confidencial"], 9, True, GRAY_HEX, PP_ALIGN.CENTER)

def _slide_conteudo(prs, titulo, corpo, tabela=None, colunas=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_texto(s, 0.6, 0.4, 8.8, 0.8, titulo, 24, True, GOLD_HEX)
    y = 1.4
    if isinstance(corpo, list):
        for p in corpo:
            _slide_texto(s, 0.8, y, 8.4, 0.6, p, 12)
            y += 0.55
    elif corpo:
        _slide_texto(s, 0.8, y, 8.4, 0.9, corpo, 12)
        y += 0.9
    if tabela and colunas:
        shape = s.shapes.add_table(len(tabela) + 1, len(colunas),
                                   Inches(0.8), Inches(y),
                                   Inches(8.4), Inches(0.4 * (len(tabela) + 1)))
        tbl = shape.table
        for j, c in enumerate(colunas):
            tbl.cell(0, j).text = c
        for i, linha in enumerate(tabela, start=1):
            for j, v in enumerate(linha):
                tbl.cell(i, j).text = str(v)

def gerar_pptx_apresentacao(lang="pt"):
    t = APRESENTACAO_TEXTOS.get(lang, APRESENTACAO_TEXTOS["pt"])
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    _slide_titulo(prs, t)
    _slide_conteudo(prs, t["sobre_t"], t["sobre_p"])
    _slide_conteudo(prs, t["mercado_t"], t["mercado_p"])
    _slide_conteudo(prs, t["port_t"], t["port_p"], t["port_linhas"], t["port_colunas"])
    _slide_conteudo(prs, t["alcance_t"], t["alcance_p"])
    _slide_conteudo(prs, t["modelo_t"], t["modelo_p"])
    _slide_conteudo(prs, t["midia_t"], t["midia_p"])
    _slide_conteudo(prs, t["midia_a_t"], "", t["midia_a_linhas"], t["midia_a_colunas"])
    _slide_conteudo(prs, t["midia_b_t"], "", t["midia_b_linhas"], t["midia_b_colunas"])
    _slide_conteudo(prs, t["b2b_t"], [t["b2b_p"]] + t["b2b_linhas"])
    _slide_conteudo(prs, t["proj_t"], [t["proj_p"]] + t["proj_linhas"])
    _slide_conteudo(prs, t["seed_t"], [t["seed_p"]] + t["seed_linhas"])
    _slide_conteudo(prs, "", [t["frase"], t["contato"], t["rodape"]])
    path = f"/tmp/Apresentacao-Slides-{lang}.pptx"
    prs.save(path)
    return path

@app.get("/api/apresentacao-slides")
def get_apresentacao_slides(lang: str = "pt"):
    if lang not in APRESENTACAO_TEXTOS:
        lang = "pt"
    path = gerar_pptx_apresentacao(lang)
    if not os.path.exists(path):
        raise HTTPException(500, "Falha ao gerar os slides")
    return FileResponse(
        path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=f"Apresentacao-Slides-{lang}.pptx")

@app.get("/api/apresentacao")
def get_apresentacao(lang: str = "pt"):
    if lang not in APRESENTACAO_TEXTOS:
        lang = "pt"
    t = APRESENTACAO_TEXTOS[lang]
    path = f"/tmp/Apresentacao-{lang}.pdf"
    doc = SimpleDocTemplate(path, pagesize=A4, leftMargin=50, rightMargin=50,
                            topMargin=50, bottomMargin=40)
    e = []
    e.append(Paragraph(t["titulo"], estilo(20, True, GOLD, TA_CENTER, 0, 4)))
    e.append(Paragraph(t["subtitulo"], estilo(12, False, GRAY, TA_CENTER, 0, 6)))
    e.append(Paragraph(t["confidencial"], estilo(8, False, GRAY, TA_CENTER, 0, 12)))
    for sec_t, sec_p in [("sobre_t","sobre_p"), ("mercado_t","mercado_p"),
                         ("port_t","port_p"), ("alcance_t","alcance_p"),
                         ("modelo_t","modelo_p"), ("midia_t","midia_p"),
                         ("b2b_t","b2b_p"), ("proj_t","proj_p"), ("seed_t","seed_p")]:
        e.append(Paragraph(t[sec_t], estilo(14, True, GOLD, TA_LEFT, 10, 4)))
        if isinstance(t[sec_p], list):
            for p in t[sec_p]:
                e.append(Paragraph(p, estilo(10, False, DARK, TA_LEFT, 0, 3)))
        else:
            e.append(Paragraph(t[sec_p], estilo(10, False, DARK, TA_LEFT, 0, 3)))
    e.append(Spacer(1, 12))
    e.append(Paragraph(t["frase"], estilo(11, True, GOLD, TA_CENTER, 10, 4)))
    e.append(Paragraph(t["contato"], estilo(9, False, GRAY, TA_CENTER, 0, 2)))
    e.append(Paragraph(t["rodape"], estilo(8, False, GRAY, TA_CENTER, 0, 2)))
    doc.build(e)
    return FileResponse(path, media_type="application/pdf", filename=f"Apresentacao-{lang}.pdf")

# ===== MODELOS PYDANTIC =====
class PayReq(BaseModel):
    nome: str
    nascimento: str
    email: Optional[str] = ""
    lang: str = "pt"
class UrnaPayReq(BaseModel):
    nome_completo: str
    nome_urna: str
    email: Optional[str] = ""
    lang: str = "pt"
    cargo: str = "vereador"
    nome1: str = ""
    nome2: str = ""
    nome3: str = ""
    nome4: str = ""
    nome5: str = ""
class EleitoralPayReq(BaseModel):
    nome_completo: str
    numero: str
    email: Optional[str] = ""
    lang: str = "pt"
    cargo: str = "vereador"
class SugestaoReq(BaseModel):
    nome: str
    email: Optional[str] = ""
    mensagem: str
class BonusReq(BaseModel):
    nome: str
    email: Optional[str] = ""
    motivo: str
class AtivarBonusReq(BaseModel):
    codigo: str
# ===== CONSTANTES DE ESTILO (PDFs) =====
GOLD = colors.HexColor("#B8860B")
LGRAY = colors.HexColor("#f0f0f0")
DARK = colors.HexColor("#222")
GRAY = colors.HexColor("#888")
FONTE = "Helvetica"
FN = "Helvetica-Bold"

def estilo(tam, bold, cor, alinh=TA_LEFT, antes=0, depois=0):
    return ParagraphStyle("estilo", fontName=FN if bold else FONTE, fontSize=tam,
                          textColor=cor, alignment=alinh, spaceBefore=antes, spaceAfter=depois)

CARGO_INFO = {
    "vereador": {"label": "Vereador"},
    "dep_estadual": {"label": "Deputado Estadual"},
    "dep_federal": {"label": "Deputado Federal"},
    "senador": {"label": "Senador"},
}
ENERGIAS = {
    1: "Lideranca", 2: "Cooperacao", 3: "Criatividade",
    4: "Trabalho", 5: "Liberdade", 6: "Familia",
    7: "Sabedoria", 8: "Poder e Prosperidade (IDEAL)", 9: "Humanitarismo",
}
# ===== GERADOR DE PDF (usa gerador_pdf.py se existir; senão fallback interno) =====
def _gerar_pdf_local(prod, data, lang, nome, bd, dado=""):
    path = f"/tmp/p_{prod}_{uuid.uuid4().hex[:8]}.pdf"
    doc = SimpleDocTemplate(path, pagesize=A4, leftMargin=50, rightMargin=50,
                            topMargin=40, bottomMargin=30)
    e = []
    titulo = PRODUTOS.get(lang, PRODUTOS["pt"]).get(prod, prod).upper()
    e.append(Spacer(1, 15))
    e.append(Paragraph(titulo, estilo(18, True, GOLD, TA_CENTER, 0, 6)))
    e.append(Paragraph(nome.upper(), estilo(12, True, DARK, TA_CENTER, 0, 2)))
    e.append(Paragraph(bd or dado, estilo(9, False, GRAY, TA_CENTER, 0, 10)))
    if isinstance(data, dict):
        td = [["Numero", "Valor"]]
        rotulos = [("life_path", "Caminho de Vida"), ("expression", "Expressao"),
                   ("soul_urge", "Motivacao"), ("personality", "Personalidade"),
                   ("destiny", "Destino")]
        for k, l in rotulos:
            if k in data:
                td.append([l, str(data[k])])
        if len(td) > 1:
            tbl = Table(td, colWidths=[200, 100])
            tbl.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), GOLD),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTSIZE", (0, 0), (-1, -1), 9),
                ("GRID", (0, 0), (-1, -1), 0.3, colors.grey),
                ("BACKGROUND", (0, 1), (-1, -1), LGRAY),
            ]))
            e.append(tbl)
    e.append(Spacer(1, 10))
    e.append(Paragraph("(c) A1ELOS", estilo(7, False, GRAY, TA_CENTER)))
    doc.build(e)
    return path
def _pagina_sucesso_local(pdf_path, nome, prod_nome, lang="pt"):
    b64 = ""
    if pdf_path and os.path.exists(pdf_path):
        with open(pdf_path, "rb") as f:
            b64 = base64.b64encode(f.read()).decode()
    btn = ""
    if b64:
        btn = (f'<a href="data:application/pdf;base64,{b64}" download="Documento.pdf" '
               f'style="display:inline-block;padding:18px 50px;background:#C9A94E;color:#000;'
               f'text-decoration:none;border-radius:50px;font-weight:700;font-size:1.2rem;margin:25px 0">📥 BAIXAR PDF</a>')
    return (f'<html><body style="background:#0a0a0a;color:#fff;text-align:center;padding:40px;'
            f'font-family:sans-serif"><h1 style="color:#C9A94E">✅ Confirmado!</h1>'
            f'<p>Ola <b>{nome}</b>, seu {prod_nome} foi gerado.</p>{btn}'
            f'<a href="/" style="color:#C9A94E">Voltar</a></body></html>')
def _entregar_arquivo_local(tipo, nome, lang="pt"):
    try:
        pf = _gerar_pdf_local(tipo, {"life_path": 1, "expression": 1, "soul_urge": 1,
                                     "personality": 1, "destiny": 1}, lang, nome, "")
        return {"pdf": pf, "url": "", "pdf_ok": True}
    except Exception as e:
        logger.error(f"Falha entrega: {e}")
        return {"pdf": None, "url": "", "pdf_ok": False}
try:
    from gerador_pdf import gerar_pdf, pagina_sucesso, _entregar_arquivo
except Exception:
    gerar_pdf = _gerar_pdf_local
    pagina_sucesso = _pagina_sucesso_local
    _entregar_arquivo = _entregar_arquivo_local
# ===== EMAIL SIMPLES =====
def _enviar_email_simples(destinatario, assunto, corpo):
    try:
        msg = MIMEMultipart()
        msg["From"] = FROM_EMAIL
        msg["To"] = destinatario
        msg["Subject"] = assunto
        msg.attach(MIMEText(corpo, "plain", "utf-8"))
        with smtplib.SMTP(os.getenv("SMTP_HOST"), int(os.getenv("SMTP_PORT", 587))) as s:
            s.starttls()
            s.login(os.getenv("SMTP_USER"), os.getenv("SMTP_PASS"))
            s.send_message(msg)
        return True
    except Exception as e:
        logger.error(f"SMTP: {e}")
        return False
# ===== CRIACAO DE SESSAO STRIPE =====
def _criar_sessao(produto, lang="pt", email="", nome="", birth="", meta_extra=None):
    if lang not in PRICE_IDS or produto not in PRICE_IDS[lang]:
        raise HTTPException(status_code=400, detail="Idioma ou produto invalido")
    price_id = PRICE_IDS[lang].get(produto, "")
    nome_prod = PRODUTOS.get(lang, PRODUTOS["pt"]).get(produto, produto)
    meta = {"tipo": produto, "lang": lang, "nome": nome, "birth": birth, "email": email}
    if meta_extra:
        meta.update(meta_extra)
    pay_types = ["card", "boleto"] if MOEDA.get(lang, "brl") == "brl" else ["card"]
    locale = lang if lang in ["pt", "en", "es", "fr", "de", "it", "ja", "zh"] else "auto"
    if produto == "urna":
        success_url = f"{BASE_URL}/api/pay/urna-success?session_id={{CHECKOUT_SESSION_ID}}"
    elif produto == "eleitoral":
        success_url = f"{BASE_URL}/api/pay/eleitoral-success?session_id={{CHECKOUT_SESSION_ID}}"
    else:
        success_url = f"{BASE_URL}/api/pay/success?session_id={{CHECKOUT_SESSION_ID}}"
    try:
        if price_id and price_id.startswith("price_"):
            session = stripe.checkout.Session.create(
                mode="payment", payment_method_types=pay_types,
                line_items=[{"price": price_id, "quantity": 1}],
                customer_email=email or None,
                locale=locale, metadata=meta,
                payment_method_options={"card": {"installments": {"enabled": True}}} if MOEDA.get(lang, "brl") == "brl" else None,
                success_url=success_url,
                cancel_url=f"{BASE_URL}/api/pay/cancel")
        else:
            session = stripe.checkout.Session.create(
                mode="payment", payment_method_types=pay_types,
                line_items=[{"price_data": {"currency": MOEDA.get(lang, "brl"),
                    "product_data": {"name": nome_prod},
                    "unit_amount": preco_local(produto, lang)}, "quantity": 1}],
                customer_email=email or None,
                locale=locale, metadata=meta,
                payment_method_options={"card": {"installments": {"enabled": True}}} if MOEDA.get(lang, "brl") == "brl" else None,
                success_url=success_url,
                cancel_url=f"{BASE_URL}/api/pay/cancel")
        return {"id": session.id, "url": session.url}
    except Exception as e:
        logger.error(f"Stripe: {e}")
        raise HTTPException(500, "Erro ao criar pagamento")
        
# ===== ROTA GENERICA /pay/{produto} =====
_ALIAS_PRODUTO = {"complete": "completo"}
@app.post("/pay/{produto}")
def pay_produto(produto: str, req: PayReq):
    if not STRIPE_KEY:
        raise HTTPException(503, "Stripe nao configurado")
    produto = _ALIAS_PRODUTO.get(produto, produto)
    lang = req.lang or "pt"
    return _criar_sessao(produto, lang, req.email, req.nome, req.nascimento)
# ===== CHECKOUT NOME DE URNA =====
@app.post("/pay/urna")
def pay_urna(req: UrnaPayReq):
    if not STRIPE_KEY:
        raise HTTPException(503, "Stripe nao configurado")
    if len(req.nome_completo.strip()) < 3:
        raise HTTPException(400, "Nome obrigatorio")
    nomes = [n.strip() for n in [req.nome1, req.nome2, req.nome3, req.nome4, req.nome5] if n.strip()]
    if not nomes:
        raise HTTPException(400, "Pelo menos 1 nome")
    meta = {"tipo": "urna", "lang": req.lang or "pt", "nome_completo": req.nome_completo,
            "cargo": req.cargo, "email": req.email, "nome": req.nome_completo}
    for i, n in enumerate(nomes, 1):
        meta[f"nome{i}"] = n
    return _criar_sessao("urna", req.lang or "pt", req.email, req.nome_completo, "", meta)
# ===== CHECKOUT NUMERO ELEITORAL =====
@app.post("/pay/eleitoral")
def pay_eleitoral(req: EleitoralPayReq):
    if not STRIPE_KEY:
        raise HTTPException(503, "Stripe nao configurado")
    if not req.numero or len(req.numero) < 2:
        raise HTTPException(400, "Numero obrigatorio")
    meta = {"tipo": "eleitoral", "lang": req.lang or "pt", "sigla": req.numero,
            "cargo": req.cargo, "email": req.email, "numero_existente": "",
            "nome_completo": req.nome_completo}
    return _criar_sessao("eleitoral", req.lang or "pt", req.email, req.nome_completo, "", meta)
# ===== CHECKOUT COLETIVO (desconto progressivo) =====
@app.get("/criar-checkout-coletivo")
async def criar_checkout_coletivo(lang: str = "pt", items: str = "[]"):
    if not STRIPE_KEY:
        raise HTTPException(503, "Stripe nao configurado")
    try:
        itens = json.loads(items)
    except Exception:
        raise HTTPException(400, "items invalidos")
    if not itens:
        raise HTTPException(400, "Nenhum item")
    qtd_total = sum(int(it.get("qtd", 0)) for it in itens if it.get("qtd"))
    desc = desconto_bc(qtd_total)
    line_items = []
    for it in itens:
        pid = it.get("id")
        qtd = int(it.get("qtd", 0))
        if not pid or qtd <= 0 or pid == "coletivo":
            continue
        unit = preco_local(pid, lang)
        if desc > 0:
            unit = int(round(unit * (1 - desc)))
        line_items.append({"price_data": {"currency": MOEDA.get(lang, "brl"),
            "product_data": {"name": PRODUTOS.get(lang, PRODUTOS["pt"]).get(pid, pid)},
            "unit_amount": unit}, "quantity": qtd})
    if not line_items:
        raise HTTPException(400, "Itens invalidos")
    pay_types = ["card", "boleto"] if MOEDA.get(lang, "brl") == "brl" else ["card"]
    locale = lang if lang in ["pt", "en", "es", "fr", "de", "it", "ja", "zh"] else "auto"
    session = stripe.checkout.Session.create(
        mode="payment", payment_method_types=pay_types,
        line_items=line_items,
        locale=locale,
        metadata={"tipo": "coletivo", "lang": lang, "desconto": str(int(desc * 100)),
                  "itens": json.dumps(itens)},
        success_url=f"{BASE_URL}/api/pay/success?session_id={{CHECKOUT_SESSION_ID}}",
        cancel_url=f"{BASE_URL}/api/pay/cancel")
    return RedirectResponse(url=session.url)
# ===== ROTA /criar-checkout (usada pelo site) — INDENTACAO CORRIGIDA =====
@app.get("/criar-checkout")
async def criar_checkout_direto(lang: str = "pt", produto: str = "express",
                                qtd: int = 0, total: float = 0, itens: str = "",
                                nome: str = "", nascimento: str = "",
                                nome_completo: str = "", cargo: str = "vereador",
                                numero: str = "", email: str = "",
                                nome1: str = "", nome2: str = "", nome3: str = "",
                                nome4: str = "", nome5: str = "",
                                energia: str = "", dado: str = "", tipo: str = ""):
    if not STRIPE_KEY:
        raise HTTPException(503, "Stripe nao configurado")
    if produto == "coletivo":
        return await criar_checkout_coletivo(lang=lang, items=itens or "[]")
    if produto not in PRODUTO_FAIXA:
        raise HTTPException(400, "Produto invalido")
    meta = {}
    if produto == "urna":
        meta = {"nome_completo": nome_completo, "cargo": cargo, "nome": nome_completo,
                "nome1": nome1, "nome2": nome2, "nome3": nome3,
                "nome4": nome4, "nome5": nome5}
    elif produto == "eleitoral":
        numero_existente: str = "",   # ← novo parâmetro
        meta = {"sigla": numero, "cargo": cargo,
            "nome_completo": nome_completo, "numero_existente": numero_existente}
    else:
        meta = {"energia": energia, "dado": dado, "tipo": tipo}
    s = _criar_sessao(produto, lang, email, nome, nascimento, meta)
    return RedirectResponse(url=s["url"])
# ===== SUCESSO POS-PAGAMENTO =====
# Produtos que coletam "dado" (nome digitado) em vez de nome+nascimento
DADO_PRODUTOS = {"nome_pet", "nickname", "nome_dominio", "nome_canal",
                 "nome_equipe", "nome_ong", "nome_projeto", "nome_evento"}
@app.get("/api/pay/success")
def pay_success(request: Request):
    sid = request.query_params.get("session_id", "")
    if not sid:
        return HTMLResponse("ERRO")
    try:
        s = stripe.checkout.Session.retrieve(sid)
        meta = getattr(s, "metadata", {}) or {}
        if hasattr(meta, "to_dict"):
            meta = meta.to_dict()
        nome = meta.get("nome", "Cliente")             
        email = meta.get("email", "") or getattr(s, "customer_email", "") or ""
        bd = meta.get("birth", "")
        prod = meta.get("tipo", "express")
        lang = meta.get("lang", "pt")
        dado = meta.get("dado", "")
        if not bd:
            bd = "2000-01-01"
        if dado:
            data = {"dado": dado}
            nome_exib = dado or nome
        elif prod in DADO_PRODUTOS:
            data = {"dado": dado or nome}
            nome_exib = dado or nome
        else:
            data = calc_mapa(nome, bd)
            nome_exib = nome
        db = SessionLocal()
        try:
            db.add(Order(id=uuid.uuid4().hex[:12], email=email or "sem-email",
                         product=prod, price=float(getattr(s, "amount_total", 0) or 0) / 100,
                         status="paid", payment_id=sid))
            db.commit()
        except Exception:
            pass
        finally:
            db.close()
        pn = PRODUTOS.get(lang, PRODUTOS["pt"]).get(prod, prod)
        pf = gerar_pdf(prod, data, lang, nome_exib, bd, dado=dado)
        html = pagina_sucesso(pf, nome_exib, pn, lang)
        if pf and os.path.exists(pf):
            os.remove(pf)
        return HTMLResponse(html)
    except Exception as e:
        logger.error(f"Success: {e}")
        return HTMLResponse("ERRO")
@app.get("/api/pay/urna-success")
def pay_urna_success(request: Request):
    sid = request.query_params.get("session_id", "")
    if not sid:
        return HTMLResponse("ERRO")
    try:
        s = stripe.checkout.Session.retrieve(sid)
        meta = getattr(s, "metadata", {}) or {}
        if hasattr(meta, "to_dict"):
            meta = meta.to_dict()
        nc = meta.get("nome_completo", "")
        cr = meta.get("cargo", "vereador")
        nomes = [meta.get(f"nome{i}", "") for i in range(1, 6) if meta.get(f"nome{i}", "")]
        if not nomes:
            return HTMLResponse("ERRO")
        res, _, sugs = validar_nomes_urna(nomes, cr)
        cl = CARGO_INFO.get(cr, {}).get("label", cr)
        lang = meta.get("lang", "pt")
        dados_urna = {"nome_completo": nc, "cargo_label": cl,
                      "resultados": res, "sugestoes": sugs}
        pf = gerar_pdf("urna", dados_urna, lang, nc, "")
        html = pagina_sucesso(pf, nc, PRODUTOS.get(lang, PRODUTOS["pt"]).get("urna", "Urna"), lang)
        if pf and os.path.exists(pf):
            os.remove(pf)
        return HTMLResponse(html)
    except Exception:
        return HTMLResponse("ERRO")
@app.get("/api/pay/eleitoral-success")
def pay_eleitoral_success(request: Request):
    sid = request.query_params.get("session_id", "")
    if not sid:
        return HTMLResponse("ERRO")
    try:
        s = stripe.checkout.Session.retrieve(sid)
        meta = getattr(s, "metadata", {}) or {}
        if hasattr(meta, "to_dict"):
            meta = meta.to_dict()
        sg = int(meta.get("sigla", "0"))
        cr = meta.get("cargo", "vereador")
        ne_str = meta.get("numero_existente", "")
        ss = str(sg).zfill(2)
        cl_map = {"vereador": "Vereador", "dep_estadual": "Dep. Estadual",
                  "dep_federal": "Dep. Federal", "senador": "Senador"}
        cl2 = cl_map.get(cr, cr)
        sugs = gerar_numeros(sg, cr)
        ni = None
        if ne_str and len(ne_str) >= 3:
            try:
                ni = {"numero": ne_str, "energia": reduzir(sum(int(d) for d in ne_str))}
            except Exception:
                pass
        lang = meta.get("lang", "pt")
        dados_ele = {"sigla": ss, "cargo_label": cl2,
                     "sugestoes": sugs, "numero_existente": ni}
        pf = gerar_pdf("eleitoral", dados_ele, lang, f"Candidato {cl2}", "")
        html = pagina_sucesso(pf, f"Candidato {cl2}",
                              PRODUTOS.get(lang, PRODUTOS["pt"]).get("eleitoral", "Eleitoral"), lang)
        if pf and os.path.exists(pf):
            os.remove(pf)
        return HTMLResponse(html)
    except Exception:
        return HTMLResponse("ERRO")
@app.get("/api/pay/cancel")
def pay_cancel():
    return HTMLResponse("<h1>Cancelado</h1><a href='/'>Voltar</a>")
# ===== CALCULO GRATIS =====
@app.post("/calculate")
def calculate(req: PayReq):
    db = SessionLocal()
    try:
        if len(req.nome.strip()) < 2:
            raise HTTPException(400, "Nome curto")
        if not req.nascimento:
            raise HTTPException(400, "Data obrigatoria")
        res = calc_mapa(req.nome, req.nascimento)
        cid = uuid.uuid4().hex[:8]
        db.add(Calc(id=cid, name=req.nome, birth_date=req.nascimento, email=req.email or "", **res))
        db.commit()
        res["download_pdf"] = False
        return {"id": cid, **res}
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Calc: {e}")
        raise HTTPException(500, "Erro")
    finally:
        db.close()

@app.post("/calculate/urna")
def calc_urna(req: UrnaPayReq):
    nomes = [n.strip() for n in [req.nome1, req.nome2, req.nome3, req.nome4, req.nome5] if n.strip()]
    res, ideal, sugs = validar_nomes_urna(nomes, req.cargo)
    return {"resultados": res, "ideal": ideal, "sugestoes": sugs}

@app.post("/calculate/eleitoral")
def calc_eleitoral(req: EleitoralPayReq):
    sigla = int(req.numero) if req.numero.isdigit() else 0
    sugs = gerar_numeros(sigla, req.cargo)
    return {"sugestoes": sugs}

# ===== ROTAS BASE =====
@app.get("/")
def root():
    try:
        return HTMLResponse(open(os.path.join(os.path.dirname(__file__), "static", "index.html"), "r", encoding="utf-8").read())
    except Exception:
        try:
            return HTMLResponse(open(os.path.join(os.path.dirname(__file__), "index.html"), "r", encoding="utf-8").read())
        except Exception:
            return HTMLResponse("<h1>API ativa</h1>")

@app.get("/config")
def config():
    return {"stripe_pk": STRIPE_PUB}

@app.get("/api/health")
def health():
    return {"status": "ok", "stripe": bool(STRIPE_KEY)}

# ===== WEBHOOK STRIPE =====
@app.post("/stripe-webhook")
async def stripe_webhook(req: Request):
    payload = await req.body()
    sig = req.headers.get("stripe-signature", "")
    whsec = os.getenv("STRIPE_WEBHOOK_SECRET", "")
    if whsec:
        try:
            event = stripe.Webhook.construct_event(payload, sig, whsec)
        except Exception:
            raise HTTPException(400, "Assinatura invalida")
    else:
        data = json.loads(payload)
        event = {"type": data.get("type", ""), "data": data.get("data", {})}
    if event["type"] == "checkout.session.completed":
        session = event["data"]["object"]
        meta = session.get("metadata", {})
        tipo = meta.get("tipo", "express")
        if tipo == "coletivo":
            try:
                items = json.loads(meta.get("itens", "[]"))
                gerados = _gerar_codigos_para_itens(items)
                logger.info(f"Coletivo: {len(gerados)} codigos gerados")
            except Exception as e:
                logger.error(f"Erro codigos coletivo: {e}")
        nome = meta.get("nome", "Cliente")
        lang = meta.get("lang", "pt")
        logger.info(f"Pagamento confirmado: {session['id']} -> {tipo}")
        entrega = _entregar_arquivo(tipo, nome, lang)
        logger.info(f"Entrega: {entrega}")
    return {"status": "success"}

# ===== SISTEMA DE BONUS =====
ARQ_BONUS = "bonus_codes.json"

def _carregar_codigos():
    try:
        with open(ARQ_BONUS, "r", encoding="utf-8") as f:
            return json.load(f)
    except FileNotFoundError:
        return {}

def _salvar_codigos(dados):
    with open(ARQ_BONUS, "w", encoding="utf-8") as f:
        json.dump(dados, f, ensure_ascii=False, indent=2)

def _gerar_codigo_bonus():
    chars = string.ascii_uppercase + string.digits
    p1 = "".join(secrets.choice(chars) for _ in range(4))
    p2 = "".join(secrets.choice(chars) for _ in range(4))
    return f"A1-{p1}-{p2}"

def _gerar_codigos_para_itens(itens):
    codigos = _carregar_codigos()
    gerados = []
    for item in itens:
        pid = item.get("id")
        qtd = int(item.get("qtd", 0))
        for _ in range(qtd):
            cod = _gerar_codigo_bonus()
            codigos[cod] = {"produto": pid, "usado": False}
            gerados.append({"codigo": cod, "produto": pid})
    _salvar_codigos(codigos)
    return gerados

@app.post("/ativar-bonus")
async def ativar_bonus(req: AtivarBonusReq):
    codigos = _carregar_codigos()
    info = codigos.get(req.codigo)
    if not info:
        return {"ok": False, "msg": "Codigo nao encontrado"}
    if info.get("usado"):
        return {"ok": False, "msg": "Codigo ja utilizado"}
    info["usado"] = True
    info["data_uso"] = datetime.now().isoformat()
    _salvar_codigos(codigos)
    target = PRODUTO_TARGET.get(info.get("produto"), "inicio")
    return {"ok": True, "target": target, "produto": info.get("produto")}

@app.post("/gerar-codigos-coletivo")
async def gerar_codigos_coletivo(req: Request):
    corpo = await req.json()
    itens = corpo.get("itens", [])
    gerados = _gerar_codigos_para_itens(itens)
    return {"ok": True, "total": len(gerados), "codigos": gerados}

# ===== CAIXA DE SUGESTOES + BONUS =====
@app.post("/sugestao")
async def receber_sugestao(req: SugestaoReq):
    try:
        with open("sugestoes.json", "a") as f:
            f.write(json.dumps({"nome": req.nome, "mensagem": req.mensagem,
                                "data": datetime.now().isoformat()}, ensure_ascii=False) + "\n")
    except Exception as e:
        logger.error(f"Erro sugestao: {e}")
    try:
        _enviar_email_simples(ADMIN_EMAIL, "Nova sugestao/reclamacao - A1ELOS",
                              f"Sugestao de {req.nome} ({req.email}):\n\n{req.mensagem}")
    except Exception as e:
        logger.error(f"Erro email sugestao: {e}")
    return {"ok": True}

@app.post("/bonus")
async def solicitar_bonus(req: BonusReq):
    try:
        with open("bonus_solicitacoes.json", "a") as f:
            f.write(json.dumps({"nome": req.nome, "motivo": req.motivo,
                                "data": datetime.now().isoformat()}, ensure_ascii=False) + "\n")
    except Exception as e:
        logger.error(f"Erro bonus: {e}")
    try:
        _enviar_email_simples(ADMIN_EMAIL, "Solicitacao de BONUS - A1ELOS",
                              f"Cliente: {req.nome}\nEmail: {req.email}\nMotivo: {req.motivo}")
    except Exception as e:
        logger.error(f"Erro email bonus: {e}")
    return {"ok": True}

# ===== SISTEMA DE PUBLICIDADE GEOLOCALIZADA =====
ARQ_BANNERS = "banners.json"
PAIS_CONTINENTE = {
    "BR":"SA","AR":"SA","CL":"SA","CO":"SA","PE":"SA","UY":"SA","PY":"SA","BO":"SA","EC":"SA","VE":"SA",
    "US":"NA","CA":"NA","MX":"NA",
    "PT":"EU","ES":"EU","FR":"EU","DE":"EU","IT":"EU","GB":"EU","RU":"EU","NL":"EU","BE":"EU","CH":"EU","AT":"EU","IE":"EU",
    "CN":"AS","JP":"AS","IN":"AS","KR":"AS","SA":"AS","AE":"AS","IL":"AS","TR":"AS","ID":"AS","PK":"AS","BD":"AS",
    "EG":"AF","NG":"AF","ZA":"AF","KE":"AF","MA":"AF",
    "AU":"OC","NZ":"OC"
}

def _carregar_banners():
    try:
        with open(ARQ_BANNERS, "r", encoding="utf-8") as f:
            return json.load(f)
    except FileNotFoundError:
        return []

def _salvar_banners(banners):
    with open(ARQ_BANNERS, "w", encoding="utf-8") as f:
        json.dump(banners, f, ensure_ascii=False, indent=2)

@app.get("/api/banner")
async def get_banner(posicao: str = "topo", pais: str = "BR"):
    banners = _carregar_banners()
    if not banners:
        return {"ok": False, "banner": None}
    continente = PAIS_CONTINENTE.get(pais.upper(), "")
    hoje = date.today().isoformat()
    for b in banners:
        if not b.get("ativo") or b.get("posicao") != posicao:
            continue
        if b.get("tipo") == "temporario":
            if b.get("data_fim") and hoje > b["data_fim"]:
                continue
            if b.get("data_inicio") and hoje < b["data_inicio"]:
                continue
        if b.get("escopo") == "pais" and b.get("pais") == pais.upper():
            return {"ok": True, "banner": b}
        if b.get("escopo") == "continente" and b.get("continente") == continente:
            return {"ok": True, "banner": b}
        if b.get("escopo") == "mundo":
            return {"ok": True, "banner": b}
    return {"ok": False, "banner": None}

# ===== INICIALIZACAO =====
if __name__ == "__main__":
    import uvicorn
    port = int(os.getenv("PORT", "10000"))
    uvicorn.run(app, host="0.0.0.0", port=port)
